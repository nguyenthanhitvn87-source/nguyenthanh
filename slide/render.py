"""Approximate renderer for QA: draws master/layout backgrounds, slide shapes,
pictures, tables and text with the real geometry and font sizes from the file.
Reports any line whose measured width exceeds its box."""
import io
import sys
from pptx import Presentation
from pptx.util import Pt
from PIL import Image, ImageDraw, ImageFont

DPI, EMU = 110, 914400
FB = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
FR = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"

# effective sizes for placeholders that inherit from the layout
PH_SIZE = {"Title 1": 35, "Subtitle 2": 16, "Text Placeholder 3": 16,
           "Text Placeholder 2": 115, "Title 4": 30}
SECTION_TITLE = 36


def px(v):
    return v / EMU * DPI


def solid(sh):
    try:
        if sh.fill.type == 1:
            return "#" + str(sh.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def draw_text(d, tf, x, y, w, h, default_size, name, sl, warn):
    lines = []
    for para in tf.paragraphs:
        size, bold, color = default_size, False, "#1A1A1A"
        text = para.text
        for r in para.runs:
            if r.font.size:
                size = r.font.size.pt
            if r.font.bold:
                bold = True
            try:
                if r.font.color and r.font.color.rgb is not None:
                    color = "#" + str(r.font.color.rgb)
            except Exception:
                pass
        align = str(para.alignment or "")
        for seg in text.replace("\n", "\x0b").split("\x0b"):
            lines.append((seg, size, bold, color, align))

    total = sum(s * 1.2 for _, s, _, _, _ in lines) * DPI / 72
    try:
        anchor = str(tf.vertical_anchor or "")
    except Exception:
        anchor = ""
    yy = y + (h - total) / 2 if "MIDDLE" in anchor or "CENTER" in anchor else y

    for text, size, bold, color, align in lines:
        f = ImageFont.truetype(FB if bold else FR, max(1, int(size * DPI / 72)))
        # wrap
        words, cur, out = text.split(" "), "", []
        for wd in words:
            t = (cur + " " + wd).strip()
            if d.textlength(t, font=f) > w and cur:
                out.append(cur)
                cur = wd
            else:
                cur = t
        out.append(cur)
        if len(out) > 1 and any(d.textlength(o, font=f) > w for o in out):
            warn.append(f"slide{sl} {name}: word too wide -> {text[:40]!r}")
        for ln in out:
            tw = d.textlength(ln, font=f)
            xx = x + (w - tw) / 2 if "CENTER" in align else x
            d.text((xx, yy), ln, font=f, fill=color)
            yy += size * 1.2 * DPI / 72
    if yy > y + h + 2:
        warn.append(f"slide{sl} {name}: text {round((yy-y)/DPI,2)}in > box {round(h/DPI,2)}in")


def render(path, out_prefix):
    prs = Presentation(path)
    W, H = int(px(prs.slide_width)), int(px(prs.slide_height))
    warn = []
    for i, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), "#FFFFFF")
        d = ImageDraw.Draw(img)
        for src in (slide.slide_layout.slide_master, slide.slide_layout):
            for sh in src.shapes:
                c = solid(sh)
                if c and sh.width and sh.height:
                    d.rectangle([px(sh.left), px(sh.top), px(sh.left + sh.width),
                                 px(sh.top + sh.height)], fill=c)
        for sh in slide.shapes:
            if sh.left is None:
                continue
            x, y, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
            if sh.shape_type == 13:
                try:
                    ic = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
                    img.paste(ic.resize((max(1, int(w)), max(1, int(h)))),
                              (int(x), int(y)), ic.resize((max(1, int(w)), max(1, int(h)))))
                except Exception:
                    d.rectangle([x, y, x + w, y + h], outline="#888888")
                continue
            if sh.has_table:
                t = sh.table
                cy = y
                for r in t.rows:
                    cx = x
                    for ci, cell in enumerate(r.cells):
                        cw = px(t.columns[ci].width)
                        d.rectangle([cx, cy, cx + cw, cy + px(r.height)], outline="#CCCCCC")
                        draw_text(d, cell.text_frame, cx + 6, cy + 4, cw - 12,
                                  px(r.height) - 8, 18, f"table r{r}c{ci}", i, warn)
                        cx += cw
                    cy += px(r.height)
                continue
            c = solid(sh)
            if c:
                if "OVAL" in str(sh.auto_shape_type or ""):
                    d.ellipse([x, y, x + w, y + h], fill=c)
                else:
                    d.rectangle([x, y, x + w, y + h], fill=c)
            if sh.has_text_frame and sh.text.strip():
                size = PH_SIZE.get(sh.name, 14)
                if sh.name == "Title 1" and slide.slide_layout.name == "Section Header":
                    size = SECTION_TITLE
                draw_text(d, sh.text_frame, x, y, w, h, size, sh.name, i, warn)
        img.save(f"{out_prefix}-{i}.png")
    for line in warn:
        print("WARN:", line)
    print("rendered", len(prs.slides), "slides")


render(sys.argv[1], sys.argv[2])
