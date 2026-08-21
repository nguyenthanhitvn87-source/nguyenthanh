"""LPNVi deck, v2 — clearer content, livelier layout, DKSH template untouched.

skel2.pptx = the original deck with master/layouts/artwork intact:
  1 title · 2 overview · 3-7 blank Section Header copies · 8 table · 9 thanks
Slides 3-7 are stripped down to their Title + big-number placeholders and
redrawn here with the template's own palette.
"""
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

RED = RGBColor(0xBE, 0x00, 0x28)
RED2 = RGBColor(0xEF, 0x23, 0x3C)
PINK = RGBColor(0xFB, 0xE7, 0xEA)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
HAIR = RGBColor(0xE9, 0xE9, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"


# ----------------------------------------------------------------- helpers
def _template_para(tf):
    for p in tf.paragraphs:
        if p.runs:
            return p._p
    return tf.paragraphs[0]._p


def set_lines(tf, lines, size=None):
    """Replace paragraphs, keeping the original run/paragraph formatting."""
    tmpl = deepcopy(_template_para(tf))
    body = tf._txBody
    for p in body.findall(qn("a:p")):
        body.remove(p)
    run_tmpl = deepcopy(tmpl.findall(qn("a:r"))[0])
    end = tmpl.find(qn("a:endParaRPr"))
    end_tmpl = deepcopy(end) if end is not None else None
    # keep only <a:pPr>; runs and endParaRPr are re-added in schema order below
    for child in list(tmpl):
        if child.tag != qn("a:pPr"):
            tmpl.remove(child)

    for line in lines:
        para = deepcopy(tmpl)
        if line:
            for i, chunk in enumerate(line.split("\n")):
                if i:
                    para.append(para.makeelement(qn("a:br"), {}))
                r = deepcopy(run_tmpl)
                r.find(qn("a:t")).text = chunk
                if size is not None:
                    rPr = r.find(qn("a:rPr"))
                    if rPr is None:
                        rPr = r.makeelement(qn("a:rPr"), {})
                        r.insert(0, rPr)
                    rPr.set("sz", str(int(size * 100)))
                para.append(r)
        # <a:endParaRPr> must be the LAST child, or PowerPoint reads the
        # paragraph as empty and shows the placeholder prompt instead
        if end_tmpl is not None:
            para.append(deepcopy(end_tmpl))
        body.append(para)


def strip(slide, keep=("Title 1", "Text Placeholder 2")):
    for sh in list(slide.shapes):
        if sh.name not in keep:
            sh._element.getparent().remove(sh._element)


def txt(slide, x, y, w, h, runs, size=12, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, wrap=True):
    """runs: a string, or a list of (text, {overrides}) — all laid out as runs
    inside ONE paragraph, so differently coloured pieces stay on the same line.
    Use "\n" inside a piece for a line break."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if space:
        p.space_after = Pt(space)
    for item in (runs if isinstance(runs, list) else [runs]):
        text, over = item if isinstance(item, tuple) else (item, {})
        for i, chunk in enumerate(text.split("\n")):
            if i:
                p.add_run().text = "\v"
            r = p.add_run()
            r.text = chunk
            r.font.name = FONT
            r.font.size = Pt(over.get("size", size))
            r.font.bold = over.get("bold", bold)
            r.font.color.rgb = over.get("color", color)
    return box


def shape(slide, kind, x, y, w, h, fill=None, line=None, lw=1.25, radius=None):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    sp.text_frame.text = ""
    return sp


def label(slide, x, y, w, h, text, fill=RED, size=10.5, color=WHITE, radius=0.5):
    sp = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, radius=radius)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return sp


def bullet(slide, x, y, w, n, head, body):
    """Numbered red square + bold heading + grey detail."""
    label(slide, x, y + 0.03, 0.34, 0.34, str(n), fill=RED, size=13, radius=0.18)
    txt(slide, x + 0.52, y, w - 0.52, 0.32, head, size=15, bold=True, color=DARK)
    txt(slide, x + 0.52, y + 0.36, w - 0.52, 0.62, body, size=12, color=GREY)


def panel(slide, x, y, w, h, title):
    """Pink framed panel with the template's red header bar."""
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=PINK, line=RED, radius=0.03)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.14, fill=RED, radius=0.5)
    txt(slide, x, y + 0.30, w, 0.30, title, size=12.5, bold=True, color=RED,
        align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                 fill=WHITE, line=RED, lw=1.0, radius=0.04)


def rule(slide, x, y, w):
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.012, fill=HAIR)


# ------------------------------------------------------------------ content
VN_NAME = "Kem dưỡng da chống nắng ABC SPF50+ dạng gel 50ml"
EN_NAME = "ABC SUNSCREEN GEL SPF50+ 50ML"

prs = Presentation("skel2.pptx")
S = list(prs.slides)


def named(slide):
    return {sh.name: sh for sh in slide.shapes}


# --------------------------------------------------------------- slide 1
s = named(S[0])
t = s["Title 1"]
t.height = Inches(1.70)
set_lines(t.text_frame,
          ["PROCESS LONG PRODUCTION NAME WITH VIETNAMESE LANGUAGE – LPNVi"],
          size=28)
set_lines(s["Subtitle 2"].text_frame, [
    "Business Owner: Luong Ngoc Truong, HEC Sales Manager",
    "Process Owner: Tran Quang Thuan, IT CSSC Vietnam",
    "Technical Owner: Nguyen Vu Thanh, IT CSSC Vietnam",
])
set_lines(s["Text Placeholder 3"].text_frame, ["HCMC, 21 Aug 2026"])

icon = s["Picture 4"]
W = Inches(2.6)
box = S[0].shapes.add_textbox(int(icon.left + icon.width / 2 - W / 2),
                              int(icon.top + icon.height + Inches(0.10)),
                              int(W), Inches(0.70))
box.name = "SAP Label"
tf = box.text_frame
tf.word_wrap = False
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "SAP"
r.font.name = FONT
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = WHITE

# --------------------------------------------------------------- slide 2
s = named(S[1])
set_lines(s["Title 1"].text_frame, ["Tổng quan giải pháp"])
for tn, cn, bn, title, cap, badge in [
    ("TextBox 7",  "TextBox 8",  "Rounded Rectangle 9",
     "Master Data\nSetup", "CMIR · MG0156 · INSP", "MASTER DATA"),
    ("TextBox 15", "TextBox 16", "Rounded Rectangle 17",
     "Billing\nOutput", "Ghép Mat. Desc + INSP", "SAP"),
    ("TextBox 23", "TextBox 24", "Rounded Rectangle 25",
     "Chứng từ\nkho / giao nhận", "PGN · PXKDL · PXKNB", "PP#"),
    ("TextBox 31", "TextBox 32", "Rounded Rectangle 33",
     "Hóa đơn\nđiện tử", "Tách separator \"| \"", "VNPT"),
]:
    set_lines(s[tn].text_frame, [title])
    set_lines(s[cn].text_frame, [cap])
    set_lines(s[bn].text_frame, [badge])
set_lines(s["TextBox 34"].text_frame,
          ["Phạm vi: all BU nationwide  —  thứ tự ưu tiên CMIR › MG0156 › INSP."])

# --------------------------------------------------------------- slide 3
sl = S[2]
strip(sl)
n = named(sl)
set_lines(n["Title 1"].text_frame, ["Hiện trạng"])
set_lines(n["Text Placeholder 2"].text_frame, ["1"])

bullet(sl, 0.85, 2.15, 5.5, 1, "Mat. Desc giới hạn 40 ký tự",
       "Tên sản phẩm dài hơn 40 ký tự bị cắt cụt khi in ra chứng từ.")
bullet(sl, 0.85, 3.40, 5.5, 2, "Đã enhance được một phần",
       "CMIR cho khách Tender, Basis Text cho nhóm hàng MG0156.")
bullet(sl, 0.85, 4.65, 5.5, 3, "Phần còn lại vẫn chưa xử lý",
       "Các BU khác chưa in được tên tiếng Việt có dấu, đầy đủ.")

panel(sl, 7.05, 1.95, 5.6, 4.55, "AS-IS  ·  Tên bị cắt")
card(sl, 7.30, 2.72, 5.10, 1.05)
txt(sl, 7.50, 2.88, 4.70, 0.22, "Tên đầy đủ cần in", size=10, color=GREY)
txt(sl, 7.50, 3.14, 4.70, 0.70, VN_NAME, size=13, bold=True, color=DARK)
label(sl, 11.05, 2.84, 1.20, 0.26, "47 KÝ TỰ", fill=RED2, size=9)

card(sl, 7.30, 4.05, 5.10, 1.28)
txt(sl, 7.50, 4.21, 4.70, 0.22, "Mat. Desc chứa được", size=10, color=GREY)
txt(sl, 7.50, 4.47, 4.70, 0.30,
    [("ABC SUNSCREEN GEL SPF50+ 50ML", {}), (" …", {"color": RED2})],
    size=13, bold=True, color=DARK)
txt(sl, 7.50, 4.83, 4.70, 0.24, "không dấu  ·  phần sau bị cắt",
    size=10, color=RED2)
label(sl, 11.05, 4.17, 1.20, 0.26, "40 KÝ TỰ", fill=GREY, size=9)

txt(sl, 7.30, 5.62, 5.10, 0.30,
    "→ Hóa đơn / phiếu kho chỉ in được tên tiếng Anh, chưa đủ nghĩa.",
    size=11, bold=True, color=RED)

# --------------------------------------------------------------- slide 4
sl = S[3]
strip(sl)
n = named(sl)
set_lines(n["Title 1"].text_frame, ["Yêu cầu"])
set_lines(n["Text Placeholder 2"].text_frame, ["2"])

bullet(sl, 0.85, 2.05, 5.5, 1, "Tên dài, có dấu tiếng Việt",
       "Vượt giới hạn 40 ký tự của Mat. Desc.")
bullet(sl, 0.85, 3.15, 5.5, 2, "Trình bày 2 dòng trên hóa đơn",
       "Dòng 1 tiếng Việt, dòng 2 tiếng Anh chữ nhỏ hơn trên EIV.")
bullet(sl, 0.85, 4.25, 5.5, 3, "Mở rộng sang chứng từ kho",
       "Phiếu Giao Nhận, Phiếu Xuất Kho Đại Lý, Phiếu Xuất Kho Nội Bộ.")
bullet(sl, 0.85, 5.35, 5.5, 4, "Phạm vi áp dụng",
       "All BU nationwide.")

panel(sl, 7.05, 1.95, 5.6, 4.55, "TO-BE  ·  Hóa đơn điện tử")
card(sl, 7.30, 2.72, 5.10, 2.05)
txt(sl, 7.50, 2.90, 4.70, 0.22, "Tên hàng hóa, dịch vụ", size=10, color=GREY)
rule(sl, 7.50, 3.18, 4.70)
txt(sl, 7.50, 3.32, 4.30, 0.62, VN_NAME, size=13, bold=True, color=DARK)
txt(sl, 7.50, 4.06, 4.30, 0.26, EN_NAME, size=9.5, color=GREY)
label(sl, 11.90, 3.34, 0.30, 0.30, "1", fill=RED, size=10, radius=0.5)
label(sl, 11.90, 4.05, 0.30, 0.30, "2", fill=GREY, size=10, radius=0.5)

txt(sl, 7.30, 5.00, 5.10, 0.26,
    [("1", {"color": RED, "bold": True}),
     ("  Tiếng Việt có dấu — chữ lớn, dòng đầu", {})], size=11, color=DARK)
txt(sl, 7.30, 5.32, 5.10, 0.26,
    [("2", {"color": GREY, "bold": True}),
     ("  Tiếng Anh — chữ nhỏ hơn, dòng dưới", {})], size=11, color=DARK)
txt(sl, 7.30, 5.78, 5.10, 0.30,
    "→ Áp dụng cho EIV và cả PGN / PXKDL / PXKNB.",
    size=11, bold=True, color=RED)

# --------------------------------------------------------------- slide 5
sl = S[4]
strip(sl)
n = named(sl)
set_lines(n["Title 1"].text_frame, ["Giải pháp"])
set_lines(n["Text Placeholder 2"].text_frame, ["3"])

# two source fields
for x, name, cap, detail, fill in [
    (1.30, "Mat. Desc", "max 40 ký tự", "Tên tiếng Anh (default)", GREY),
    (5.55, "INSP", "max 220 ký tự", "Tên tiếng Việt có dấu", RED),
]:
    card(sl, x, 1.95, 3.35, 1.20)
    label(sl, x + 0.22, 2.14, 1.55, 0.30, name, fill=fill, size=11.5)
    txt(sl, x + 1.95, 2.16, 1.20, 0.26, cap, size=10.5, bold=True, color=fill,
        align=PP_ALIGN.RIGHT)
    txt(sl, x + 0.22, 2.62, 2.95, 0.50, detail, size=12, color=DARK)

txt(sl, 4.65, 2.38, 0.90, 0.50, "+", size=26, bold=True, color=RED,
    align=PP_ALIGN.CENTER)

# concatenation result
shape(sl, MSO_SHAPE.DOWN_ARROW, 4.90, 3.30, 0.40, 0.44, fill=RED)
txt(sl, 5.45, 3.40, 3.40, 0.34, "nối bằng separator  \"| \"",
    size=12, bold=True, color=RED)

card(sl, 1.30, 3.92, 7.60, 0.95)
txt(sl, 1.52, 4.08, 7.16, 0.24, "Chuỗi gửi sang EIV VNPT", size=10, color=GREY)
txt(sl, 1.52, 4.36, 7.16, 0.44,
    [("Mat. Desc", {"color": GREY}),
     ("  |  ", {"color": RED2, "bold": True}),
     ("INSP text", {"color": RED})],
    size=15, bold=True)

# priority chain
txt(sl, 1.30, 5.25, 3.20, 0.28, "Thứ tự ưu tiên logic", size=11.5, bold=True,
    color=DARK)
chain = [("CMIR", "khách Tender"), ("MG0156", "Basis Text"), ("INSP", "phần còn lại")]
cx = 1.30
for i, (nm, why) in enumerate(chain):
    card(sl, cx, 5.55, 2.20, 0.80)
    txt(sl, cx, 5.69, 2.20, 0.26, nm, size=13, bold=True, color=RED,
        align=PP_ALIGN.CENTER)
    txt(sl, cx, 5.97, 2.20, 0.24, why, size=9.5, color=GREY, align=PP_ALIGN.CENTER)
    if i < 2:
        shape(sl, MSO_SHAPE.CHEVRON, cx + 2.32, 5.79, 0.32, 0.32, fill=RED)
    cx += 2.70

panel(sl, 9.45, 1.95, 3.25, 4.40, "KEY POINT")
txt(sl, 9.70, 2.68, 2.75, 0.85,
    "INSP chứa được 220 ký tự và giữ nguyên dấu tiếng Việt.",
    size=13, bold=True, color=DARK)
rule(sl, 9.70, 3.62, 2.75)
txt(sl, 9.70, 3.78, 2.75, 0.80,
    "Mat. Desc giữ nguyên như hiện tại — không phải sửa master data cũ.",
    size=11, color=GREY)
rule(sl, 9.70, 4.70, 2.75)
txt(sl, 9.70, 4.86, 2.75, 0.85,
    "VNPT tách chuỗi theo \"| \" để xếp thành 2 dòng trên hóa đơn.",
    size=11, color=GREY)

# --------------------------------------------------------------- slide 6
sl = S[5]
strip(sl)
n = named(sl)
set_lines(n["Title 1"].text_frame, ["Ví dụ thực tế"])
set_lines(n["Text Placeholder 2"].text_frame, ["4"])

# input fields
txt(sl, 0.85, 1.95, 11.85, 0.28, "Master data nhập vào", size=11.5, bold=True,
    color=DARK)
card(sl, 0.85, 2.30, 5.75, 0.95)
txt(sl, 1.05, 2.44, 5.35, 0.22, "Mat. Desc", size=10, bold=True, color=GREY)
txt(sl, 1.05, 2.72, 5.35, 0.30, EN_NAME, size=12, color=DARK)
card(sl, 6.95, 2.30, 5.75, 0.95)
txt(sl, 7.15, 2.44, 5.35, 0.22, "INSP", size=10, bold=True, color=RED)
txt(sl, 7.15, 2.72, 5.35, 0.30, VN_NAME, size=12, color=DARK)

# concatenated string
txt(sl, 0.85, 3.55, 11.85, 0.28, "Chuỗi SAP ghép và gửi sang EIV",
    size=11.5, bold=True, color=DARK)
card(sl, 0.85, 3.90, 11.85, 0.72)
txt(sl, 1.05, 4.08, 11.45, 0.36,
    [(EN_NAME, {"color": GREY}),
     ("|  ", {"color": RED2}),
     (VN_NAME, {"color": RED})],
    size=12.5, bold=True)

txt(sl, 1.05, 4.68, 11.45, 0.24,
    [("\u25cf Mat. Desc — tiếng Anh", {"color": GREY}),
     ("      \u25cf separator \"| \"", {"color": RED2}),
     ("      \u25cf INSP — tiếng Việt có dấu", {"color": RED})],
    size=9.5, bold=True)

# rendered invoice
txt(sl, 0.85, 5.02, 11.85, 0.28, "Kết quả trình bày trên hóa đơn / phiếu kho",
    size=11.5, bold=True, color=DARK)
panel(sl, 0.85, 5.36, 11.85, 1.66, "")
card(sl, 1.10, 5.56, 11.35, 1.24)
txt(sl, 1.35, 5.70, 10.85, 0.22, "Tên hàng hóa, dịch vụ", size=9.5, color=GREY)
rule(sl, 1.35, 5.96, 10.85)
txt(sl, 1.35, 6.08, 10.85, 0.32, VN_NAME, size=14, bold=True, color=DARK)
txt(sl, 1.35, 6.46, 10.85, 0.26, EN_NAME, size=10, color=GREY)

# --------------------------------------------------------------- slide 7
sl = S[6]
strip(sl)
n = named(sl)
set_lines(n["Title 1"].text_frame, ["Lập trình"])
set_lines(n["Text Placeholder 2"].text_frame, ["5"])

streams = [
    ("SAP", "Billing Output",
     ["Ghép Mat. Desc + INSP text", "Chèn separator \"| \"",
      "Áp thứ tự CMIR › MG0156 › INSP"], RED),
    ("PP#", "Form chứng từ kho",
     ["Phiếu Giao Nhận (PGN)", "Phiếu Xuất Kho Đại Lý (PXKDL)",
      "Phiếu Xuất Kho Nội Bộ (PXKNB)"], RED),
    ("VNPT", "Hóa đơn điện tử EIV",
     ["Nhận biết separator \"| \"", "Tách thành 2 dòng",
      "Tiếng Anh cỡ chữ nhỏ hơn"], GREY),
]
x = 0.85
for owner, title, items, tone in streams:
    card(sl, x, 2.00, 3.85, 4.35)
    shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.00, 3.85, 0.13, fill=tone, radius=0.5)
    label(sl, x + 0.28, 2.36, 1.25, 0.34, owner, fill=tone, size=12)
    txt(sl, x + 0.28, 2.92, 3.30, 0.56, title, size=15, bold=True, color=DARK)
    rule(sl, x + 0.28, 3.60, 3.30)
    yy = 3.80
    for it in items:
        label(sl, x + 0.28, yy + 0.04, 0.13, 0.13, "", fill=tone, radius=0.5)
        txt(sl, x + 0.56, yy, 3.02, 0.62, it, size=11.5, color=DARK)
        yy += 0.72
    x += 4.16

txt(sl, 0.85, 6.60, 11.85, 0.30,
    "Ba nhánh chạy song song — master data setup xong trước khi go-live.",
    size=11.5, color=GREY)

# --------------------------------------------------------------- slide 8
s = named(S[7])
set_lines(s["Title 1"].text_frame, ["Tóm lại"])
set_lines(s["TextBox 2"].text_frame, ["CMIR › MG0156 › INSP"])
table = next(sh for sh in S[7].shapes if sh.has_table).table
for r, values in enumerate([
    ["#", "Nội dung", "Ai làm", "Trạng thái"],
    ["1", "Phạm vi: all BU nationwide", "Business", "Đã chốt"],
    ["2", "Master data: CMIR, MG0156, INSP", "Master Data", "Setup"],
    ["3", "Ưu tiên: CMIR › MG0156 › INSP", "SAP", "Lập trình"],
    ["4", "Dentsply: xóa hết record CMIR trước", "CIT", "Lưu ý"],
]):
    for c, value in enumerate(values):
        set_lines(table.cell(r, c).text_frame, [value])

prs.save("lpnvi.pptx")
print("saved lpnvi.pptx")
