"""Add 'SAP' directly under the icon on the title slide of the original deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation("src.pptx")
slide = prs.slides[0]

icon = next(sh for sh in slide.shapes if sh.name == "Picture 4")
cx = icon.left + icon.width / 2
top = icon.top + icon.height + Inches(0.10)

W = Inches(2.6)
box = slide.shapes.add_textbox(int(cx - W / 2), int(top), int(W), Inches(0.55))
box.name = "SAP Label"

tf = box.text_frame
tf.word_wrap = False
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
tf.vertical_anchor = MSO_ANCHOR.TOP

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "SAP"
run.font.name = "Arial"
run.font.size = Pt(28)
run.font.bold = True
run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

prs.save("out.pptx")
print("icon box (in):", icon.left / 914400, icon.top / 914400, icon.width / 914400, icon.height / 914400)
print("SAP box  (in):", box.left / 914400, box.top / 914400, box.width / 914400, box.height / 914400)
